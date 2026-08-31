import glob, json, re
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

CHUNK_SIZE = 800      # 조각 하나 최대 글자수
OVERLAP = 100         # 조각끼리 겹치는 글자수

def read_pdf(path):
    reader = PdfReader(path)
    return "\n".join((p.extract_text() or "") for p in reader.pages)

def read_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def read_xlsx(path):
    wb = load_workbook(path, data_only=True)
    lines = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)

def read_pptx(path):
    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)

def split_text(text, size=CHUNK_SIZE, overlap=OVERLAP):
    text = re.sub(r"\s+", " ", text).strip()  # 공백 정리
    chunks = []
    start = 0
    while start < len(text):
        chunks.append(text[start:start + size])
        start += size - overlap
    return chunks

all_chunks = []

# ---- 1) 연금 문서 (pdf/docx/xlsx/pptx) ----
for path in sorted(glob.glob("docs/pension_pdf/docs_renamed/*")):
    try:
        if path.endswith(".pdf"):
            text = read_pdf(path)
        elif path.endswith(".docx"):
            text = read_docx(path)
        elif path.endswith(".xlsx"):
            text = read_xlsx(path)
        elif path.endswith(".pptx"):
            text = read_pptx(path)
        else:
            continue
        if len(text.strip()) < 50:
            print(f"[스킵-스캔본] {path}")
            continue
        for i, c in enumerate(split_text(text)):
            all_chunks.append({"source": path.split("/")[-1],
                               "type": "pension", "chunk_id": i, "text": c})
        print(f"[완료] {path.split('/')[-1]}")
    except Exception as e:
        print(f"[에러] {path}: {e}")

# ---- 2) 투자설명서 (펀드코드별 폴더) ----
for path in sorted(glob.glob("docs/fund_pdf/*/*/*.pdf")):
    try:
        fund_code = path.split("/")[-2]   # 폴더명 = 펀드 코드
        text = read_pdf(path)
        if len(text.strip()) < 50:
            print(f"[스킵-스캔본] {path}")
            continue
        for i, c in enumerate(split_text(text)):
            all_chunks.append({"source": path.split("/")[-1],
                               "type": "fund", "fund_code": fund_code,
                               "chunk_id": i, "text": c})
        print(f"[완료] {fund_code}")
    except Exception as e:
        print(f"[에러] {path}: {e}")

with open("chunks.json", "w", encoding="utf-8") as f:
    json.dump(all_chunks, f, ensure_ascii=False, indent=1)

print(f"\n총 {len(all_chunks)}개 청크 저장 완료 → chunks.json")
